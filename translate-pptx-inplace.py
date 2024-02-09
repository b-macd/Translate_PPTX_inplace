from pptx import Presentation
from transformers import pipeline



def translate_pptx_inplace(input_file, output_file):
    
    # Instantiate translation pipeline
    pipe = pipeline("translation", model="Helsinki-NLP/opus-mt-ar-en")

    # Load the presentation
    prs = Presentation(input_file)

    # Helsinki model has a bug that replaces a empty prompt with this obnoxious string.
    filler = 'Hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey.'
    
    # For each slide in the presentation
    for slide_number, slide in enumerate(prs.slides):  
        
        # For each shape in a slide
        for shape in slide.shapes:
    
            # Testing the "has_text_frame" parameter
            if shape.has_text_frame:
                
                # For each paragraph of the text_frame
                for paragraph in shape.text_frame.paragraphs:
                    
                    # Send paragraphs through translation pipeline
                    results = pipe(paragraph.text)
                    text_to_add = results[0]['translation_text']
                    text_fixed = text_to_add.replace(filler, '')
                    paragraph.text = text_fixed
                    
    # save new file with translations
    return prs.save(output_file)

if __name__ == "__main__":
    
    file_name = './Presentation1.pptx'
    new_file = 'test2.pptx'
    texts = translate_pptx_inplace(file_name, new_file)