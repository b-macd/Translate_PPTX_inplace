import streamlit as st # pip install streamlit==0.82.0
from pptx import Presentation
from transformers import pipeline
from io import BytesIO


st.set_page_config(page_title='Powerpoint Translator', layout='wide', initial_sidebar_state='expanded')


st.title('Powerpoint Translator')


uploaded_file = st.file_uploader('upload your powerpoint file here')

if uploaded_file:
   filename = st.write("Filename: ", uploaded_file.name)
   out_name = uploaded_file.name.replace('.pptx', '')

if st.button('Translate powerpoint'):
        
    # Instantiate translation pipeline
    pipe = pipeline("translation", model="Helsinki-NLP/opus-mt-ar-en")

    # Load the presentation
    prs = Presentation(uploaded_file)

    # Helsinki model has a bug that replaces a empty prompt with this obnoxious string.
    filler = 'Hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey.'
    
    # For each slide in the presentation
    for slide_number, slide in enumerate(prs.slides):  
        
        # For each shape in a slide
        for shape in slide.shapes:
    
            # Testing the "has_text_frame" parameter
            if shape.has_text_frame:
                
                # For each paragraph of the text_frame
                for paragraph in shape.text_frame.paragraphs:
                    
                    # Send paragraphs through translation pipeline
                    results = pipe(paragraph.text)
                    text_to_add = results[0]['translation_text']
                    text_fixed = text_to_add.replace(filler, '')
                    paragraph.text = text_fixed
    binary_output = BytesIO()
# save new file with translations
    prs.save(binary_output)


    st.download_button(label='Click to download PowerPoint',data=binary_output.getvalue(),file_name=f'{out_name}-translated.pptx')



 
 