import streamlit as st
from pptx import Presentation
from transformers import MarianMTModel, AutoTokenizer
from io import BytesIO


st.set_page_config(page_title='Powerpoint Translator', layout='wide', initial_sidebar_state='expanded')


st.title('Powerpoint Translator')


uploaded_file = st.file_uploader('upload your powerpoint file here')

if uploaded_file:
   filename = st.write("Filename: ", uploaded_file.name)
   out_name = uploaded_file.name.replace('.pptx', '')

if st.button('Translate powerpoint'):
        
    # Instantiate translation pipeline
    def translation_pipeline(original_text):
        model_name = 'Helsinki-NLP/opus-mt-ar-en'
        model = MarianMTModel.from_pretrained(model_name)
        tokenizer = AutoTokenizer.from_pretrained(model_name)
        batch = tokenizer([original_text], return_tensors= 'pt')
        generated_ids = model.generate(**batch)
        translated_text = tokenizer.batch_decode(generated_ids, skipped_special_tokens=True)[0]
        return translated_text

    # Load the presentation
    prs = Presentation(uploaded_file)

    # Helsinki model has a bug that replaces a empty prompt with this obnoxious string.
    filler = 'Hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey.'
    
    # For each slide in the presentation
    for slide_number, slide in enumerate(prs.slides):  
        
        # For each shape in a slide
        for shape in slide.shapes:
    
            # Testing the "has_text_frame" parameter
            if shape.has_text_frame:
                
                # For each paragraph of the text_frame
                for paragraph in shape.text_frame.paragraphs:
                    
                    # Send paragraphs through translation pipeline
                    results = translation_pipeline(paragraph.text)
                    text_fixed = results.replace(filler, '')
                    paragraph.text = text_fixed
            
            elif shape.has_table:
                table_filler = 'Hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey, hey.'
                tbl = shape.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                for r in range(0, row_count):
                    for c in range(0, col_count):
                        cell = tbl.cell(r, c)
                        paragraphs = cell.text_frame.paragraphs
                        for paragraph in paragraphs:
                            results = translation_pipeline(paragraph.text)
                            text_fixed = results.replace(filler, '')
                            paragraph.text = text_fixed

    binary_output = BytesIO()
# save new file with translations
    prs.save(binary_output)


    st.download_button(label='Click to download PowerPoint',data=binary_output.getvalue(),file_name=f'{out_name}-translated.pptx')
